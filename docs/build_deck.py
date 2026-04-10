"""
DataPrismAI — Demo Presentation Generator
Run: python docs/build_deck.py
Output: docs/DataPrismAI_Demo.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand colours ──────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy background
TEAL       = RGBColor(0x00, 0xC9, 0xB1)   # primary accent
GOLD       = RGBColor(0xF5, 0xA6, 0x23)   # secondary accent
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xD6, 0xE0)
DARK_GRAY  = RGBColor(0x44, 0x55, 0x66)
GREEN      = RGBColor(0x2E, 0xCC, 0x71)
RED        = RGBColor(0xE7, 0x4C, 0x3C)

W  = Inches(13.33)   # widescreen width
H  = Inches(7.5)     # widescreen height


# ── Helpers ────────────────────────────────────────────────────────────────────
def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape


def add_text(slide, text, x, y, w, h,
             size=20, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, size=16, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, space_before=Pt(4)):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def accent_bar(slide, color=TEAL, height=Inches(0.06)):
    add_rect(slide, 0, 0, W, height, color)


def bottom_bar(slide, color=TEAL, height=Inches(0.06)):
    add_rect(slide, 0, H - height, W, height, color)


def slide_number(slide, num: int, total: int):
    add_text(slide, f"{num} / {total}",
             W - Inches(1.2), H - Inches(0.4), Inches(1.0), Inches(0.3),
             size=10, color=DARK_GRAY, align=PP_ALIGN.RIGHT)


def section_chip(slide, label: str, x=Inches(0.35), y=Inches(0.15)):
    """Small coloured pill label in top-left."""
    chip = add_rect(slide, x, y, Inches(1.6), Inches(0.28), TEAL)
    add_text(slide, label.upper(),
             x + Inches(0.08), y + Inches(0.02), Inches(1.5), Inches(0.26),
             size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


# ── Slide builders ─────────────────────────────────────────────────────────────

def slide_cover(prs):
    s = blank_slide(prs)
    fill_bg(s, NAVY)

    # Left accent stripe
    add_rect(s, 0, 0, Inches(0.5), H, TEAL)

    # Gradient-ish right panel
    add_rect(s, Inches(8.2), 0, Inches(5.13), H, RGBColor(0x06, 0x24, 0x3C))

    # Logo / name
    add_text(s, "DataPrism", Inches(0.8), Inches(1.6), Inches(5), Inches(1.0),
             size=52, bold=True, color=WHITE)
    add_text(s, "AI", Inches(0.8) + Inches(3.65), Inches(1.6), Inches(2), Inches(1.0),
             size=52, bold=True, color=TEAL)

    add_text(s, "Conversational Analytics for Enterprise Banking",
             Inches(0.8), Inches(2.65), Inches(7), Inches(0.6),
             size=20, color=LIGHT_GRAY)

    add_text(s, "GenBI  ·  Semantic SQL  ·  3-Mode AI  ·  Real-time Insights",
             Inches(0.8), Inches(3.2), Inches(7.5), Inches(0.5),
             size=14, color=DARK_GRAY, italic=True)

    # Tags
    for i, tag in enumerate(["Pattern Mode", "Hybrid Mode", "LLM Mode"]):
        tx = Inches(0.8) + i * Inches(1.85)
        add_rect(s, tx, Inches(3.9), Inches(1.6), Inches(0.34), DARK_GRAY)
        add_text(s, tag, tx + Inches(0.08), Inches(3.92), Inches(1.5), Inches(0.3),
                 size=11, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # Right panel content
    add_text(s, "DEMO", Inches(8.6), Inches(1.4), Inches(4), Inches(0.8),
             size=48, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    add_text(s, "April 2026", Inches(8.6), Inches(2.3), Inches(4), Inches(0.5),
             size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    stats = [
        ("64", "Tables"),
        ("600K+", "Deposit Txns"),
        ("3", "Countries"),
        ("3", "Chat Modes"),
    ]
    for i, (val, lbl) in enumerate(stats):
        ry = Inches(3.2) + i * Inches(0.85)
        add_rect(s, Inches(8.8), ry, Inches(3.6), Inches(0.72), RGBColor(0x10, 0x2A, 0x40))
        add_text(s, val, Inches(8.9), ry + Inches(0.04), Inches(1.2), Inches(0.64),
                 size=26, bold=True, color=TEAL)
        add_text(s, lbl, Inches(10.1), ry + Inches(0.22), Inches(2.1), Inches(0.3),
                 size=14, color=LIGHT_GRAY)

    bottom_bar(s, TEAL)


def slide_agenda(prs, num, total):
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    accent_bar(s)
    bottom_bar(s)
    slide_number(s, num, total)

    add_text(s, "Agenda", Inches(0.6), Inches(0.4), Inches(6), Inches(0.7),
             size=34, bold=True, color=WHITE)
    add_rect(s, Inches(0.6), Inches(1.15), Inches(0.5), Inches(0.05), TEAL)

    items = [
        ("01", "Problem & Opportunity",    "Why enterprise analytics is broken"),
        ("02", "What is DataPrismAI?",     "Architecture & value proposition"),
        ("03", "Data Platform",            "64-table banking data model, 3 countries"),
        ("04", "3 Chat Modes",             "Pattern → Hybrid → LLM"),
        ("05", "Semantic Layer",           "Metrics, dimensions, data dictionary"),
        ("06", "Live Demo Walkthrough",    "Real queries, real data"),
        ("07", "Personas",                 "Analyst · CFO · Fraud Analyst · Manager"),
        ("08", "Architecture Deep-dive",   "FastAPI · StarRocks · Ollama · Next.js"),
        ("09", "Roadmap",                  "What's next"),
    ]

    for i, (num_s, title, sub) in enumerate(items):
        row_y = Inches(1.35) + i * Inches(0.63)
        add_rect(s, Inches(0.6), row_y, Inches(0.5), Inches(0.44),
                 TEAL if i == 0 else RGBColor(0x10, 0x2A, 0x40))
        add_text(s, num_s, Inches(0.62), row_y + Inches(0.04), Inches(0.46), Inches(0.38),
                 size=13, bold=True, color=NAVY if i == 0 else TEAL, align=PP_ALIGN.CENTER)
        add_text(s, title, Inches(1.25), row_y + Inches(0.04), Inches(5.5), Inches(0.22),
                 size=14, bold=True, color=WHITE)
        add_text(s, sub, Inches(1.25), row_y + Inches(0.26), Inches(5.5), Inches(0.2),
                 size=10, color=DARK_GRAY, italic=True)


def _section_divider(prs, label, subtitle, num, total):
    s = blank_slide(prs)
    fill_bg(s, RGBColor(0x06, 0x14, 0x20))
    add_rect(s, 0, 0, Inches(0.5), H, TEAL)
    add_rect(s, Inches(0.5), H / 2 - Inches(0.03), W, Inches(0.06), RGBColor(0x0D, 0x1B, 0x2A))
    add_text(s, label, Inches(0.9), H / 2 - Inches(1.1), Inches(10), Inches(1.0),
             size=44, bold=True, color=WHITE)
    add_text(s, subtitle, Inches(0.9), H / 2 - Inches(0.15), Inches(10), Inches(0.5),
             size=18, color=TEAL, italic=False)
    bottom_bar(s)
    slide_number(s, num, total)
    return s


def slide_problem(prs, num, total):
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    accent_bar(s)
    bottom_bar(s)
    slide_number(s, num, total)
    section_chip(s, "Problem")

    add_text(s, "Enterprise Analytics is Stuck in the Past",
             Inches(0.5), Inches(0.5), Inches(12), Inches(0.7),
             size=30, bold=True, color=WHITE)

    pains = [
        ("📊", "Data locked in BI tools",      "Analysts wait days for dashboards. Questions can't be answered without IT."),
        ("🔍", "SQL gatekeeping",               "Business users can't query data directly — every insight needs an engineer."),
        ("🧩", "No semantic layer",             "The same metric calculated 5 different ways. No single source of truth."),
        ("🤖", "LLM hallucinations",            "Generic AI tools generate wrong SQL against unknown schemas."),
        ("⏱",  "Time to insight is too slow",  "Reports take hours. Decisions are made on stale data."),
    ]

    for i, (icon, title, body) in enumerate(pains):
        col = i % 3
        row = i // 3
        cx = Inches(0.5) + col * Inches(4.2)
        cy = Inches(1.35) + row * Inches(2.1)
        w  = Inches(3.9)
        card = add_rect(s, cx, cy, w, Inches(1.85), RGBColor(0x10, 0x2A, 0x40))
        add_text(s, icon, cx + Inches(0.15), cy + Inches(0.12), Inches(0.5), Inches(0.5), size=22)
        add_text(s, title, cx + Inches(0.65), cy + Inches(0.15), w - Inches(0.8), Inches(0.38),
                 size=14, bold=True, color=WHITE)
        add_text(s, body, cx + Inches(0.15), cy + Inches(0.6), w - Inches(0.3), Inches(1.1),
                 size=11, color=LIGHT_GRAY)


def slide_what_is(prs, num, total):
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    accent_bar(s)
    bottom_bar(s)
    slide_number(s, num, total)
    section_chip(s, "Solution")

    add_text(s, "What is DataPrismAI?",
             Inches(0.5), Inches(0.5), Inches(10), Inches(0.7),
             size=30, bold=True, color=WHITE)

    add_text(s,
             "An enterprise-grade conversational analytics platform that turns natural language "
             "into governed SQL — with semantic context, persona-aware insights, and multi-mode AI.",
             Inches(0.5), Inches(1.2), Inches(8.5), Inches(0.8),
             size=15, color=LIGHT_GRAY)

    pillars = [
        (TEAL,  "💬",  "Natural Language\nto SQL",      "Ask in plain English.\nGet accurate, governed SQL instantly."),
        (GOLD,  "🧠",  "Semantic Layer",                 "Metrics, dimensions, and joins\ndefined once — reused everywhere."),
        (GREEN, "🎭",  "Persona-Aware\nInsights",        "CFO, Analyst, Fraud Analyst\neach get tailored responses."),
        (RGBColor(0x9B, 0x59, 0xB6), "⚡", "3-Mode AI\nEngine",  "Pattern · Hybrid · LLM\n— pick speed or intelligence."),
    ]

    for i, (color, icon, title, body) in enumerate(pillars):
        cx = Inches(0.5) + i * Inches(3.1)
        cy = Inches(2.15)
        w  = Inches(2.8)
        add_rect(s, cx, cy, w, Inches(4.5), RGBColor(0x10, 0x2A, 0x40))
        add_rect(s, cx, cy, w, Inches(0.06), color)
        add_text(s, icon,  cx + Inches(0.15), cy + Inches(0.2), Inches(0.5), Inches(0.5), size=26)
        add_text(s, title, cx + Inches(0.15), cy + Inches(0.75), w - Inches(0.3), Inches(0.7),
                 size=14, bold=True, color=WHITE)
        add_text(s, body,  cx + Inches(0.15), cy + Inches(1.5), w - Inches(0.3), Inches(2.8),
                 size=11, color=LIGHT_GRAY)


def slide_data_platform(prs, num, total):
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    accent_bar(s)
    bottom_bar(s)
    slide_number(s, num, total)
    section_chip(s, "Data Platform")

    add_text(s, "Banking Data Platform — 3 Countries, 64 Tables",
             Inches(0.5), Inches(0.5), Inches(12), Inches(0.7),
             size=30, bold=True, color=WHITE)

    layers = [
        (TEAL,  "Semantic Layer",   "semantic_*",   "12 tables",  "Customer 360 · Portfolio KPIs · Spend Metrics · Risk Metrics · Payment Status"),
        (GOLD,  "Data Products",    "dp_*",         "8 tables",   "Aggregated analytical outputs ready for charting and reporting"),
        (GREEN, "Domain Data Model","ddm_*",        "14 tables",  "Conformed & enriched — standardised business entities"),
        (RGBColor(0x9B,0x59,0xB6), "Raw Layer", "raw_*", "18 tables", "Deposit txns · Deposit accounts · Loans · Repayments · Customers"),
    ]

    for i, (color, label, prefix, count, desc) in enumerate(layers):
        cy = Inches(1.35) + i * Inches(1.35)
        # arrow funnel effect
        indent = i * Inches(0.3)
        w = W - Inches(1.0) - indent * 2
        add_rect(s, Inches(0.5) + indent, cy, w, Inches(1.1), RGBColor(0x10, 0x2A, 0x40))
        add_rect(s, Inches(0.5) + indent, cy, Inches(0.08), Inches(1.1), color)
        add_text(s, label, Inches(0.75) + indent, cy + Inches(0.08), Inches(2.4), Inches(0.38),
                 size=14, bold=True, color=WHITE)
        add_text(s, f"{prefix}  |  {count}", Inches(0.75) + indent, cy + Inches(0.52),
                 Inches(2.4), Inches(0.26), size=11, bold=False, color=color)
        add_text(s, desc, Inches(3.4) + indent, cy + Inches(0.28), Inches(8.5) - indent * 2,
                 Inches(0.56), size=11, color=LIGHT_GRAY)

    # Stats bar
    stats = [("SG · MY · IN", "Countries"), ("600K+", "Deposit Transactions"),
             ("4,112", "Loans"), ("136K", "Loan Repayments")]
    bary = Inches(6.8)
    for i, (val, lbl) in enumerate(stats):
        bx = Inches(0.5) + i * Inches(3.0)
        add_text(s, val, bx, bary, Inches(2.5), Inches(0.35), size=18, bold=True, color=TEAL)
        add_text(s, lbl, bx, bary + Inches(0.35), Inches(2.5), Inches(0.25),
                 size=10, color=LIGHT_GRAY)


def slide_3_modes(prs, num, total):
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    accent_bar(s)
    bottom_bar(s)
    slide_number(s, num, total)
    section_chip(s, "AI Engine")

    add_text(s, "3 Chat Modes — Speed or Intelligence, You Choose",
             Inches(0.5), Inches(0.5), Inches(12), Inches(0.7),
             size=30, bold=True, color=WHITE)

    modes = [
        (TEAL,  "⚡",  "Pattern",  "Deterministic SQL only\n— no LLM involved",
         ["Instant response (<50ms)", "100% predictable output", "Pre-built analytical patterns",
          "Direct table lookups", "Best for: operational dashboards"]),
        (GOLD,  "★",  "Hybrid",   "Pattern first → LLM fallback\n— Recommended",
         ["Pattern matched first", "Falls back to Qwen 2.5 32B", "Full semantic context",
          "Data dictionary enrichment", "Best for: most business queries"]),
        (GREEN, "🤖", "LLM",      "Pure LLM with full\ndata dictionary context",
         ["Skips pattern matching", "Full DDL + dict injected", "Most flexible reasoning",
          "Handles complex joins", "Best for: exploratory / ad-hoc"]),
    ]

    for i, (color, icon, title, sub, bullets) in enumerate(modes):
        cx = Inches(0.45) + i * Inches(4.25)
        cy = Inches(1.35)
        w  = Inches(3.95)
        h  = Inches(5.6)
        card = add_rect(s, cx, cy, w, h, RGBColor(0x0A, 0x20, 0x30))
        add_rect(s, cx, cy, w, Inches(0.08), color)

        add_text(s, icon,  cx + Inches(0.2),  cy + Inches(0.2),  Inches(0.6), Inches(0.55), size=28)
        add_text(s, title, cx + Inches(0.85), cy + Inches(0.25), w - Inches(1.0), Inches(0.4),
                 size=22, bold=True, color=color)
        add_text(s, sub,   cx + Inches(0.2),  cy + Inches(0.8),  w - Inches(0.4), Inches(0.65),
                 size=12, color=LIGHT_GRAY, italic=True)

        for j, bullet in enumerate(bullets):
            by = cy + Inches(1.6) + j * Inches(0.68)
            add_rect(s, cx + Inches(0.2), by + Inches(0.12), Inches(0.12), Inches(0.12), color)
            add_text(s, bullet, cx + Inches(0.45), by, w - Inches(0.65), Inches(0.55),
                     size=12, color=WHITE)

        if i == 1:  # badge "recommended"
            add_rect(s, cx + w - Inches(1.5), cy + Inches(0.15), Inches(1.35), Inches(0.28),
                     GOLD)
            add_text(s, "RECOMMENDED",
                     cx + w - Inches(1.48), cy + Inches(0.17), Inches(1.3), Inches(0.25),
                     size=8, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


def slide_semantic_layer(prs, num, total):
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    accent_bar(s)
    bottom_bar(s)
    slide_number(s, num, total)
    section_chip(s, "Semantic Layer")

    add_text(s, "Governed Semantic Layer — One Source of Truth",
             Inches(0.5), Inches(0.5), Inches(12), Inches(0.7),
             size=30, bold=True, color=WHITE)

    # Left: semantic tables
    tables = [
        ("semantic_customer_360",     "600 rows",  "Full customer snapshot: balance, utilisation, overdue status"),
        ("semantic_deposit_portfolio","600 rows",  "Savings · Current · FD balances per customer per month"),
        ("semantic_loan_portfolio",   "511 rows",  "Personal · Home · Auto loan positions + NPL flags"),
        ("semantic_spend_metrics",    "9,125 rows","Monthly spend by category per customer"),
        ("semantic_portfolio_kpis",   "Monthly",   "Country-level KPIs: delinquency, utilisation, fraud rate"),
        ("semantic_payment_status",   "Per acct",  "Payment status, overdue days, amount still owed"),
        ("semantic_risk_metrics",     "Monthly",   "Fraud rate, decline rate, avg fraud score by segment"),
    ]

    for i, (tname, count, desc) in enumerate(tables):
        cy = Inches(1.3) + i * Inches(0.78)
        add_rect(s, Inches(0.5), cy, Inches(7.5), Inches(0.66), RGBColor(0x0A, 0x20, 0x30))
        add_text(s, tname, Inches(0.65), cy + Inches(0.06), Inches(3.2), Inches(0.28),
                 size=11, bold=True, color=TEAL)
        add_text(s, count, Inches(3.9), cy + Inches(0.06), Inches(1.1), Inches(0.28),
                 size=11, bold=False, color=GOLD)
        add_text(s, desc, Inches(5.1), cy + Inches(0.06), Inches(2.9), Inches(0.5),
                 size=10, color=LIGHT_GRAY)

    # Right: dictionary tables
    add_text(s, "Data Dictionary", Inches(8.2), Inches(1.25), Inches(4.5), Inches(0.38),
             size=16, bold=True, color=TEAL)
    dict_items = [
        ("dic_tables",        "Table definitions + layer + domain"),
        ("dic_columns",       "Column descriptions + business rules"),
        ("dic_relationships", "Join paths between tables (FK model)"),
        ("semantic_metrics",  "Metric definitions (calculated KPIs)"),
        ("semantic_dimensions","Dimension catalogue"),
        ("semantic_glossary_metrics", "Business glossary terms"),
    ]
    for i, (tname, desc) in enumerate(dict_items):
        cy = Inches(1.75) + i * Inches(0.78)
        add_rect(s, Inches(8.2), cy, Inches(4.7), Inches(0.64), RGBColor(0x06, 0x24, 0x3C))
        add_text(s, tname, Inches(8.35), cy + Inches(0.06), Inches(2.4), Inches(0.26),
                 size=10, bold=True, color=GOLD)
        add_text(s, desc, Inches(8.35), cy + Inches(0.33), Inches(4.4), Inches(0.26),
                 size=10, color=LIGHT_GRAY)


def slide_demo_questions(prs, num, total):
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    accent_bar(s)
    bottom_bar(s)
    slide_number(s, num, total)
    section_chip(s, "Live Demo")

    add_text(s, "Demo Query Playbook",
             Inches(0.5), Inches(0.5), Inches(12), Inches(0.7),
             size=30, bold=True, color=WHITE)

    cols = [
        ("⚡ Pattern Mode", TEAL, [
            "How many customers?",
            "Show payment status breakdown",
            "Show portfolio KPIs",
            "Top spending customers",
            "Spend by category",
            "Show monthly transaction trend",
        ]),
        ("★ Hybrid Mode", GOLD, [
            "How many customers by country?",
            "Show deposit portfolio summary",
            "NPL rate by country",
            "Show overdue customers",
            "Show loan portfolio overview",
            "Delinquency rate breakdown",
        ]),
        ("🤖 LLM Mode", GREEN, [
            "Which customers have the highest\nutilization and are overdue?",
            "Compare loan NPL rates across\ncountries for HNW segment",
            "Show FD maturity risk next 30 days\nby country",
        ]),
    ]

    for i, (title, color, questions) in enumerate(cols):
        cx = Inches(0.4) + i * Inches(4.2)
        cy = Inches(1.35)
        w  = Inches(3.9)
        add_rect(s, cx, cy, w, Inches(0.42), color)
        add_text(s, title, cx + Inches(0.15), cy + Inches(0.06), w - Inches(0.3), Inches(0.3),
                 size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

        for j, q in enumerate(questions):
            qy = cy + Inches(0.55) + j * Inches(0.88)
            add_rect(s, cx, qy, w, Inches(0.78), RGBColor(0x0A, 0x20, 0x30))
            add_text(s, f'"{q}"', cx + Inches(0.15), qy + Inches(0.1), w - Inches(0.3), Inches(0.6),
                     size=11, color=WHITE, italic=True)


def slide_personas(prs, num, total):
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    accent_bar(s)
    bottom_bar(s)
    slide_number(s, num, total)
    section_chip(s, "Personas")

    add_text(s, "Persona-Aware AI — Different Lenses, Same Data",
             Inches(0.5), Inches(0.5), Inches(12), Inches(0.7),
             size=30, bold=True, color=WHITE)

    personas = [
        ("📊", "Analyst",       TEAL,
         "Detailed breakdowns,\nSQL shown,\nnext-step suggestions"),
        ("💼", "CFO",           GOLD,
         "P&L framing, risk exposure,\nregulatory implications,\nhigh-level KPIs"),
        ("🕵", "Fraud Analyst", RED,
         "Suspicious patterns, fraud scores,\ndecline reasons,\nhigh-risk accounts"),
        ("👔", "Manager",       GREEN,
         "Team performance, targets,\noperational summary,\nno technical jargon"),
        ("🌏", "Regional Risk", RGBColor(0x9B,0x59,0xB6),
         "Country-level risk,\nNPL by region,\ndelinquency buckets"),
        ("🛍", "Retail User",   RGBColor(0xE6,0x7E,0x22),
         "Account balance, spend summary,\npayment due dates,\nsimple language"),
    ]

    for i, (icon, name, color, desc) in enumerate(personas):
        col = i % 3
        row = i // 3
        cx = Inches(0.5)  + col * Inches(4.15)
        cy = Inches(1.4)  + row * Inches(2.6)
        w, h = Inches(3.85), Inches(2.3)
        add_rect(s, cx, cy, w, h, RGBColor(0x0A, 0x20, 0x30))
        add_rect(s, cx, cy, w, Inches(0.07), color)
        add_text(s, icon, cx + Inches(0.2), cy + Inches(0.15), Inches(0.55), Inches(0.55), size=28)
        add_text(s, name, cx + Inches(0.8), cy + Inches(0.18), w - Inches(1.0), Inches(0.4),
                 size=18, bold=True, color=color)
        add_text(s, desc, cx + Inches(0.2), cy + Inches(0.8), w - Inches(0.4), Inches(1.3),
                 size=12, color=LIGHT_GRAY)


def slide_architecture(prs, num, total):
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    accent_bar(s)
    bottom_bar(s)
    slide_number(s, num, total)
    section_chip(s, "Architecture")

    add_text(s, "Architecture — End-to-End Stack",
             Inches(0.5), Inches(0.5), Inches(12), Inches(0.7),
             size=30, bold=True, color=WHITE)

    layers = [
        (TEAL,  "Presentation",  "Next.js + Tailwind",
         "Chat · Explorer · Reports · Audit · Settings\n3 chat modes UI · Persona selector · Thread history"),
        (GOLD,  "API & AI",      "FastAPI + LangGraph",
         "Graph-based query pipeline: Planner → Semantic Resolver → SQL Node → Execution → Response\nOllama · Qwen2.5:32b · Vanna · SchemaRegistry (live DESCRIBE)"),
        (GREEN, "Semantic & SQL","StarRocks 4.0 + SchemaRegistry",
         "64-table cc_analytics DB · 4-layer medallion · Live schema discovery · Auto-refresh every 5 min\nPattern engine · Canonical SQL catalog · Dimensional routing"),
        (RGBColor(0x9B,0x59,0xB6), "Data & Infra", "Docker Compose + Postgres + Ollama + MinIO",
         "StarRocks FE/BE · PostgreSQL (metadata) · Ollama (LLM inference) · MinIO (object store)\nPersistent volumes · Health-check polling · Idempotent setup script"),
    ]

    for i, (color, layer, tech, desc) in enumerate(layers):
        cy = Inches(1.3) + i * Inches(1.45)
        add_rect(s, Inches(0.5), cy, Inches(12.3), Inches(1.3), RGBColor(0x0A, 0x20, 0x30))
        add_rect(s, Inches(0.5), cy, Inches(0.08), Inches(1.3), color)
        add_text(s, layer, Inches(0.75), cy + Inches(0.1), Inches(1.6), Inches(0.35),
                 size=13, bold=True, color=color)
        add_text(s, tech,  Inches(0.75), cy + Inches(0.5), Inches(1.6), Inches(0.35),
                 size=10, color=LIGHT_GRAY, italic=True)
        add_text(s, desc,  Inches(2.5), cy + Inches(0.15), Inches(10.1), Inches(0.98),
                 size=10, color=WHITE)


def slide_how_it_works(prs, num, total):
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    accent_bar(s)
    bottom_bar(s)
    slide_number(s, num, total)
    section_chip(s, "How It Works")

    add_text(s, "How a Query Flows Through DataPrismAI",
             Inches(0.5), Inches(0.5), Inches(12), Inches(0.7),
             size=30, bold=True, color=WHITE)

    steps = [
        (TEAL,  "1", "User Message",        "Natural language question\nentered in chat UI"),
        (GOLD,  "2", "Planner Node",         "Intent classification:\nmetric_query · preview_data · schema_query"),
        (GREEN, "3", "Semantic Resolver",    "Matches metric + dimensions\nfrom Semantic Catalog"),
        (TEAL,  "4", "SQL Node",             "Pattern → Analytical → Canonical\n→ LLM fallback → Vanna"),
        (GOLD,  "5", "StarRocks Execution",  "Runs SQL against cc_analytics\nwith timeout + row cap"),
        (GREEN, "6", "Response Node",        "Persona-aware narrative + KPIs\n+ chart type + follow-ups"),
    ]

    box_w = Inches(1.85)
    arrow_w = Inches(0.3)
    total_w = len(steps) * box_w + (len(steps) - 1) * arrow_w
    start_x = (W - total_w) / 2

    for i, (color, num_s, title, desc) in enumerate(steps):
        cx = start_x + i * (box_w + arrow_w)
        cy = Inches(1.6)
        add_rect(s, cx, cy, box_w, Inches(3.6), RGBColor(0x0A, 0x20, 0x30))
        add_rect(s, cx, cy, box_w, Inches(0.07), color)
        add_rect(s, cx + Inches(0.1), cy + Inches(0.15), Inches(0.38), Inches(0.38),
                 color)
        add_text(s, num_s, cx + Inches(0.1), cy + Inches(0.15), Inches(0.38), Inches(0.38),
                 size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(s, title, cx + Inches(0.1), cy + Inches(0.65), box_w - Inches(0.2),
                 Inches(0.55), size=12, bold=True, color=WHITE)
        add_text(s, desc,  cx + Inches(0.1), cy + Inches(1.25), box_w - Inches(0.2),
                 Inches(2.1), size=10, color=LIGHT_GRAY)

        # Arrow between boxes
        if i < len(steps) - 1:
            ax = cx + box_w + Inches(0.05)
            ay = cy + Inches(1.7)
            add_text(s, "→", ax, ay, arrow_w, Inches(0.4), size=18, color=TEAL,
                     align=PP_ALIGN.CENTER)


def slide_roadmap(prs, num, total):
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    accent_bar(s)
    bottom_bar(s)
    slide_number(s, num, total)
    section_chip(s, "Roadmap")

    add_text(s, "What's Next",
             Inches(0.5), Inches(0.5), Inches(12), Inches(0.7),
             size=30, bold=True, color=WHITE)

    now = [
        "✅  Natural language to SQL (Pattern + Hybrid + LLM)",
        "✅  64-table banking data model (3 countries)",
        "✅  Semantic catalog with governed metrics & dimensions",
        "✅  6 persona types with tailored response styles",
        "✅  SchemaRegistry — live auto-refreshing schema discovery",
        "✅  Data dictionary (dic_tables / dic_columns / dic_relationships)",
        "✅  Persistent Docker volumes — data survives restarts",
    ]
    next_items = [
        "🔜  Chart rendering with auto-recommended chart types",
        "🔜  Trino federation for cross-source queries",
        "🔜  Metric Explorer UI with drill-down navigation",
        "🔜  Skills Marketplace (pluggable AI skills)",
        "🔜  Inline rich cards (KPI tiles, mini-charts in chat)",
        "🔜  Audit trail & query lineage panel",
        "🔜  Multi-thread conversation memory",
    ]

    add_text(s, "Delivered", Inches(0.5), Inches(1.3), Inches(6), Inches(0.35),
             size=15, bold=True, color=TEAL)
    for i, item in enumerate(now):
        cy = Inches(1.75) + i * Inches(0.66)
        add_text(s, item, Inches(0.5), cy, Inches(5.8), Inches(0.55),
                 size=12, color=WHITE)

    add_text(s, "Coming Next", Inches(6.8), Inches(1.3), Inches(6), Inches(0.35),
             size=15, bold=True, color=GOLD)
    for i, item in enumerate(next_items):
        cy = Inches(1.75) + i * Inches(0.66)
        add_text(s, item, Inches(6.8), cy, Inches(5.8), Inches(0.55),
                 size=12, color=LIGHT_GRAY)

    # Divider
    add_rect(s, Inches(6.5), Inches(1.3), Inches(0.04), Inches(5.5), DARK_GRAY)


def slide_closing(prs, num, total):
    s = blank_slide(prs)
    fill_bg(s, RGBColor(0x06, 0x14, 0x20))
    add_rect(s, 0, 0, Inches(0.6), H, TEAL)

    add_text(s, "DataPrism", Inches(1.0), Inches(2.0), Inches(8), Inches(1.1),
             size=56, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(s, "AI", Inches(1.0) + Inches(6.1), Inches(2.0), Inches(3), Inches(1.1),
             size=56, bold=True, color=TEAL, align=PP_ALIGN.LEFT)

    add_text(s, "Ask your data anything.",
             Inches(1.0), Inches(3.15), Inches(10), Inches(0.6),
             size=24, color=LIGHT_GRAY)

    add_text(s, "⚡ Pattern  ·  ★ Hybrid  ·  🤖 LLM",
             Inches(1.0), Inches(3.85), Inches(10), Inches(0.45),
             size=15, bold=True, color=TEAL)

    bottom_bar(s, TEAL)
    slide_number(s, num, total)


# ── Main ───────────────────────────────────────────────────────────────────────

def build():
    prs = new_prs()

    slides_fn = [
        ("cover",           slide_cover),
        ("agenda",          slide_agenda),
        ("problem",         slide_problem),
        ("what_is",         slide_what_is),
        ("data_platform",   slide_data_platform),
        ("3_modes",         slide_3_modes),
        ("semantic",        slide_semantic_layer),
        ("demo_questions",  slide_demo_questions),
        ("personas",        slide_personas),
        ("how_it_works",    slide_how_it_works),
        ("architecture",    slide_architecture),
        ("roadmap",         slide_roadmap),
        ("closing",         slide_closing),
    ]

    TOTAL = len(slides_fn)

    for idx, (name, fn) in enumerate(slides_fn, 1):
        if name == "cover":
            fn(prs)
        elif name == "closing":
            fn(prs, idx, TOTAL)
        else:
            fn(prs, idx, TOTAL)

    out = "/home/acs1980/workspace/dataprismai/docs/DataPrismAI_Demo.pptx"
    prs.save(out)
    print(f"Saved → {out}")


if __name__ == "__main__":
    build()
